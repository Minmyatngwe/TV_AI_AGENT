from prompts import summarize_prompt,selecting_image_prompt,defining_layout_prompt,layout_variant_prompt,slide_critique_prompt,html_fix_prompt
import requests 
from bs4 import BeautifulSoup
import re
from langchain_ollama import ChatOllama
from langchain_core.prompts import ChatPromptTemplate
from langchain_core.prompts import MessagesPlaceholder
import os 
from langchain_core.messages import AIMessage,HumanMessage
import base64
import qrcode
import time
import json
import playwright 
from playwright.sync_api import sync_playwright
import ollama 
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE 
import ast 


summarize_template=ChatPromptTemplate(
    [
        ("system",summarize_prompt),
        ("human","{input}")
    ]
)

ollama_model_text="qwen3.5:9b"
ollama_model_vl="qwen3-vl:9b"
deepseek_r1="hengwen/DeepSeek-R1-Distill-Qwen-32B:q4_k_m"

qwen_35="qwen3.5:35b"

summarize_model=ChatOllama(model=qwen_35)
summarize_chain=summarize_template|summarize_model





import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

def extract_from_shapes(shapes, placeholder_list):
    """
    Recursive function to find text and image placeholders 
    in a collection of shapes (handles nested groups).
    """
    
    for shape in shapes:
        if "image" in shape.name.lower():
            placeholder_list.append(shape.name)

        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            extract_from_shapes(shape.shapes, placeholder_list)

        elif shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                clean_text = paragraph.text.strip()
                if clean_text.startswith("{{") and clean_text.endswith("}}"):
                    full_placeholder = ""
                    for i,run in enumerate(paragraph.runs):
                        if i<1:
                            full_placeholder=run.text
                        else:
                            raise ValueError("Only one style is allow per text shape")
                        
                    placeholder_list.append(full_placeholder)
                            
                                                

def get_placeholder_name(slides_path):
    print("Processing slides...")
    slide_placeholder = {}
    
    for path in slides_path:
        slide_file_name = os.path.basename(path)
        slide_placeholder[slide_file_name] = []
        
        print(f"Reading: {slide_file_name}")
        presentation = Presentation(path)
        
        for slide in presentation.slides:
            extract_from_shapes(slide.shapes, slide_placeholder[slide_file_name])
            
    for key in slide_placeholder:
        slide_placeholder[key] = list(set(slide_placeholder[key]))
            
    print(slide_placeholder)
    return slide_placeholder


def recursive_replace(shapes,ai_response,slide_name,folder_path):
    for slide_dict in ai_response:
        if list(slide_dict.keys())[0]==slide_name:
            tempo_dict=slide_dict[slide_name]
    for shape in shapes:
        if "image" in shape.name.lower():
            image_path=os.path.join(folder_path,tempo_dict[shape.name.lower()])
            left = shape.left
            top = shape.top
            width = shape.width
            height = shape.height
            sp = shape._element
            sp.getparent().remove(sp)

            shapes.add_picture(image_path, left, top, width=width, height=height)
            
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            recursive_replace(shape.shapes,ai_response,slide_name,folder_path)
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                clean_text = paragraph.text.strip()
                if clean_text.startswith("{{") and clean_text.endswith("}}"):
                    paragraph.text=tempo_dict[clean_text]
    

            

def replace_placeholder_text(slides_path,ai_response,file_path):
    powerpoint_file_path=[]
    for i,path in enumerate(slides_path):
        presentation=Presentation(path)
        slide_name=os.path.basename(path)
        for slide in presentation.slides:
            recursive_replace(slide.shapes,ai_response,slide_name,file_path)
        powerpoint_save_path_parent=os.path.join(file_path,"powerpoint")
        os.makedirs(powerpoint_save_path_parent,exist_ok=True)
        full_powerpoint_file_path=os.path.join(powerpoint_save_path_parent,f"slide_{i}.pptx")
        
        presentation.save(full_powerpoint_file_path)
        powerpoint_file_path.append(full_powerpoint_file_path)
    return powerpoint_file_path
        
        
        
        
            
            
def summarize(link,slides_path):
        global summarize_chain
        response=requests.get(link)
        soup=BeautifulSoup(response.content,'html.parser')
        text=soup.text
        text_clean=re.sub(r"\n+","\n",text)
        
        images=soup.find_all('img',class_="fl-photo-img")

        slide_placeholders=get_placeholder_name(slides_path)
        
        summarization=summarize_chain.invoke({"input":text_clean,"slide_placeholder":slide_placeholders})
        print(summarization.content)
        try:
            summarization=json.loads(summarization.content)
            folder_name=summarization[-1]
            
            os.makedirs("./output",exist_ok=True)
            full_folder_path=os.path.join("./output",folder_name)
            os.makedirs(full_folder_path,exist_ok=True)
            print(f"THe folder is saved to the {full_folder_path}")

            for j,i in enumerate(images):
                im=i.get("data-src")
                if im==None:
                    im=i.get("src")
                re_image=requests.get(im)

                download_path=f"download_image{j}.png"
                print(im)
                with open(os.path.join(full_folder_path,download_path),"wb") as d:
                    d.write(re_image.content)
            files_names=[]
            multi_images=[]
            for i in os.listdir(full_folder_path):
                if i.lower().endswith(('.png', '.jpg', '.jpeg', '.webp')):
                        full_path = os.path.join(full_folder_path, i)
                        multi_images.append(full_path)
                        files_names.append(i)

            filenames_text = "\n".join(f"{i+1}. {fn}" for i, fn in enumerate(files_names))
            
            json_summarization=json.dumps(summarization[:-1])
            print(json_summarization)
            agent_2_response=ollama.chat(
                model=qwen_35,
                messages=[{
                    "role":"system","content":selecting_image_prompt
                },
                {
                    "role":"user","content":f"Topic: {json_summarization}\n\nSelect the best image. Return ONLY the filename. Here is file name {filenames_text}", "images": multi_images
                }]
            )
            qr_image=qrcode.make(link)
            
            qr_path=os.path.join(full_folder_path,"qr_code.png")
            
            qr_image.save(qr_path)

            return agent_2_response['message']['content'],full_folder_path
        except Exception as e:
            return e

def generate_template(link,slides_path):
    
    try:
        raw_ai_text, file_path = summarize(link, slides_path)
        print(raw_ai_text)
        start_idx = raw_ai_text.find('[')
        end_idx = raw_ai_text.rfind(']') + 1
        clean_list_string = raw_ai_text[start_idx:end_idx]

        ai_response = json.loads(clean_list_string)

        print("Success! AI Data:", ai_response)
        print("Folder Path:", file_path)        
        print(ai_response)
        powerpoint_paths=replace_placeholder_text(slides_path,ai_response,file_path)
        return file_path,powerpoint_paths

    except Exception as e :
        return e
if __name__=="__main__":
    print(generate_template("https://www.roboai.fi/en/news-en/apply-to-study-at-the-roboai-academy-2026/",["./template/testing_.pptx","./template/second_testing_template.pptx"]))
