from prompts import summarize_prompt,selecting_image_prompt
import requests 
from bs4 import BeautifulSoup
import re
from langchain_ollama import ChatOllama
from langchain_core.prompts import ChatPromptTemplate
from langchain_core.prompts import MessagesPlaceholder
import os 
from langchain_core.messages import AIMessage,HumanMessage
import base64
import qrcode
import time
import json
import playwright 
from playwright.sync_api import sync_playwright
import ollama 
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE 
import ast 
import subprocess
from pdf2image import convert_from_path
import shutil
import traceback

summarize_template=ChatPromptTemplate(
    [
        ("system",summarize_prompt),
        ("human","{input}")
    ]
)

ollama_model_text="qwen3.5:9b"
ollama_model_vl="qwen3-vl:30b"

qwen_35="qwen3.5:35b"

summarize_model=ChatOllama(model=qwen_35)
summarize_chain=summarize_template|summarize_model





import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

def extract_from_shapes(shapes, placeholder_list):
    """
    Recursive function to find text and image placeholders 
    in a collection of shapes (handles nested groups).
    """
    
    for shape in shapes:
        shape_name = shape.name.lower().strip()
        if "image" in shape_name or "qr_code" in shape_name:
            placeholder_list.append(shape.name)

        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            extract_from_shapes(shape.shapes, placeholder_list)

        elif shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                clean_text = paragraph.text.strip()
                if clean_text.startswith("{{") and clean_text.endswith("}}"):
                    full_placeholder = ""
                    for i,run in enumerate(paragraph.runs):
                        if i<1:
                            full_placeholder=run.text
                        else:
                            raise ValueError("Only one style is allow per text shape")
                        
                    placeholder_list.append(full_placeholder)
                            
                                                

def get_placeholder_name(slides_path):
    print("Processing slides...")
    slide_placeholder = {}
    
    for path in slides_path:
        slide_file_name = os.path.basename(path)
        slide_placeholder[slide_file_name] = []
        
        print(f"Reading: {slide_file_name}")
        print(path)
        presentation = Presentation(path)
        
        for slide in presentation.slides:
            extract_from_shapes(slide.shapes, slide_placeholder[slide_file_name])
            
    for key in slide_placeholder:
        slide_placeholder[key] = list(set(slide_placeholder[key]))
            
    print(slide_placeholder)
    return slide_placeholder


def normalize_placeholder(text):
    return text.replace("\xa0", " ").strip()

def recursive_replace(shapes, ai_response, slide_name, folder_path):
    tempo_dict = None
    for slide_dict in ai_response:
        if slide_name in slide_dict:
            tempo_dict = slide_dict[slide_name]
            break

    if tempo_dict is None:
        return

    # support both "title" and "{{title}}"
    normalized_map = {}
    for k, v in tempo_dict.items():
        key = normalize_placeholder(k)
        normalized_map[key] = v

        if key.startswith("{{") and key.endswith("}}"):
            stripped = key[2:-2].strip()
            normalized_map[stripped] = v
        else:
            normalized_map[f"{{{{{key}}}}}"] = v

    for shape in list(shapes):
        shape_name = normalize_placeholder(shape.name)

        # image replacement
        if "image" in shape_name.lower():
            image_value = (
                normalized_map.get(shape.name)
                or normalized_map.get(normalize_placeholder(shape.name))
                or normalized_map.get(shape.name.lower())
                or normalized_map.get("main_image")
                or normalized_map.get("{{main_image}}")
            )

            if image_value:
                image_path = os.path.join(folder_path, image_value)
                left, top, width, height = shape.left, shape.top, shape.width, shape.height
                sp = shape._element
                sp.getparent().remove(sp)
                shapes.add_picture(image_path, left, top, width=width, height=height)
            continue

        # qr replacement
        if "qr_code" in shape_name.lower():
            qr_path = os.path.join(folder_path, "qr_code.png")
            left, top, width, height = shape.left, shape.top, shape.width, shape.height
            sp = shape._element
            sp.getparent().remove(sp)
            shapes.add_picture(qr_path, left, top, width=width, height=height)
            continue

        # group
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            recursive_replace(shape.shapes, ai_response, slide_name, folder_path)
            continue

        # text replacement
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                clean_text = normalize_placeholder(paragraph.text)

                if clean_text.startswith("{{") and clean_text.endswith("}}"):
                    replacement = (
                        normalized_map.get(clean_text)
                        or normalized_map.get(clean_text[2:-2].strip())
                    )

                    if replacement is not None and len(paragraph.runs) == 1:
                        paragraph.runs[0].text = str(replacement)                        
def replace_placeholder_text(slides_path,ai_response,file_path):
    powerpoint_file_path=[]
    for i,path in enumerate(slides_path):
        print(path)
        presentation=Presentation(path)
        slide_name=os.path.basename(path)
        base_name=os.path.splitext(slide_name)[0]
        for slide in presentation.slides:
            recursive_replace(slide.shapes,ai_response,slide_name,file_path)
        powerpoint_save_path_parent=os.path.join(file_path,"powerpoint")
        os.makedirs(powerpoint_save_path_parent,exist_ok=True)
        full_powerpoint_file_path=os.path.join(powerpoint_save_path_parent,f"{base_name}.pptx")
        
        presentation.save(full_powerpoint_file_path)
        powerpoint_file_path.append(full_powerpoint_file_path)
    return powerpoint_file_path
        
        



def convert_pptx_to_png(pptx_path):
    # 1. Force the path to be absolute so LibreOffice doesn't get confused
    abs_pptx_path = os.path.abspath(pptx_path)
    
    output_path = os.path.dirname(os.path.dirname(abs_pptx_path))
    pdf_path = os.path.join(output_path, "pdf_files")
    os.makedirs(pdf_path, exist_ok=True)    
    
    base_name = os.path.splitext(os.path.basename(abs_pptx_path))[0]
    tempo_path = r"/tmp/pdf_tempo"  
    os.makedirs(tempo_path, exist_ok=True) 
    
    try:
        # 2. Add the temporary profile argument back in!
        profile_arg = f"-env:UserInstallation=file://{tempo_path}/profile"
        
        subprocess.run([
            "libreoffice", 
            profile_arg, 
            "--headless", 
            "--convert-to", "pdf", 
            "--outdir", tempo_path, 
            abs_pptx_path # Use the absolute path here
        ], check=True)
        
    except subprocess.CalledProcessError as e:
        print(f"Error rendering PDF: {e}")
        return []    
        
    temp_pdf_file = os.path.join(tempo_path, f"{base_name}.pdf")
    pdf_file_path = os.path.join(pdf_path, f"{base_name}.pdf")
    
    # 3. Safely move the file
    if os.path.exists(temp_pdf_file):
        shutil.move(temp_pdf_file, pdf_file_path)
    else:
        print("Error: The PDF was not found in the temp folder. LibreOffice failed silently.")
        return []

    try:
        images = convert_from_path(pdf_file_path, dpi=300)
    except Exception as e:
        print(f"Error converting PDF to images: {e}")
        return []
    
    image_path = os.path.join(output_path, "image")
    os.makedirs(image_path, exist_ok=True)    
    
    
    for i, image in enumerate(images):
        png_filename = f"{base_name}.png"
        png_full_path = os.path.join(image_path, png_filename)
        
        image.save(png_full_path, "PNG")
        print(f"Saved: {png_filename}")
            
    return png_full_path

    

def summarize(link,slides_path):
        global summarize_chain
        response=requests.get(link)
        soup=BeautifulSoup(response.content,'html.parser')
        text=soup.text
        text_clean=re.sub(r"\n+","\n",text)
        
        images = soup.select('div[class*="fl-photo-img-"]')

        slide_placeholders=get_placeholder_name(slides_path)
        print(slide_placeholders)
        summarization=summarize_chain.invoke({"input":text_clean,"slide_placeholder":slide_placeholders})
        print(summarization.content)
        try:
            summarization=json.loads(summarization.content)
            folder_name=summarization[-1]
            
            os.makedirs("./output",exist_ok=True)
            full_folder_path=os.path.join("./output",folder_name)
            os.makedirs(full_folder_path,exist_ok=True)
            print(f"THe folder is saved to the {full_folder_path}")
            image_counter = 0

            for j, div in enumerate(images):
                a_tag = div.find("a")
                image_tag = div.find("img")
                image_url = None

                if a_tag:
                    image_url = a_tag.get("href")
                elif image_tag:
                    image_url = image_tag.get("data-src") or image_tag.get("src")

                if not image_url or image_url.startswith("data:image/svg"):
                    continue

                if image_url.lower().endswith((".png", ".jpg", ".jpeg", ".webp")):
                    re_image = requests.get(image_url)

                    download_path = f"download_image{image_counter}.png"
                    with open(os.path.join(full_folder_path, download_path), "wb") as d:
                        d.write(re_image.content)

                    image_counter += 1
            print("images")
            print(images)
            files_names=[]
            multi_images=[]
            for i in os.listdir(full_folder_path):
                if i.lower().endswith(('.png', '.jpg', '.jpeg', '.webp')):
                        full_path = os.path.join(full_folder_path, i)
                        multi_images.append(full_path)
                        files_names.append(i)
            print("multi image")
            print(multi_images)

            filenames_text = "\n".join(f"{i+1}. {fn}" for i, fn in enumerate(files_names[:5]))
            json_summarization=json.dumps(summarization[:-1])
            
            agent_2_response=ollama.chat(
                model=qwen_35,
                messages=[{
                    "role":"system","content":selecting_image_prompt
                },
                {
                    "role":"user","content": f"Input list: {json_summarization}\n\nCandidate image file names:\n{filenames_text}\n\nReturn the full updated JSON list only.", "images": multi_images[:5]
                }]
            )
            
            qr_image=qrcode.make(link)
            
            qr_path=os.path.join(full_folder_path,"qr_code.png")
            
            qr_image.save(qr_path)

            return agent_2_response['message']['content'],full_folder_path,text_clean
        except Exception as e:
            print("ERROR IN summarize():", e)
            traceback.print_exc()
            raise
def generate_template(link,slides_path):
    
    try:
        raw_ai_text, file_path,web_text = summarize(link, slides_path)
        start_idx = raw_ai_text.find('[')
        end_idx = raw_ai_text.rfind(']') + 1
        clean_list_string = raw_ai_text[start_idx:end_idx]

        ai_response = json.loads(clean_list_string)

        print("Success! AI Data:", ai_response)
        print("Folder Path:", file_path)        
        print(ai_response)
        powerpoint_paths=replace_placeholder_text(slides_path,ai_response,file_path)
        
        png_image_path=[]
        for i in powerpoint_paths:
            png_path=convert_pptx_to_png(i)
            png_image_path.append(png_path)
        return file_path,powerpoint_paths,png_image_path,web_text,ai_response

    except Exception:
            import traceback
            traceback.print_exc()
            raise
if __name__=="__main__":
    print(generate_template("https://www.roboai.fi/en/we-research-en/we-research-gerotestbed-model-supporting-the-development-of-agetechnologies/",["./template/testing_.pptx","./template/second_testing_template.pptx","./template/template1.pptx"]))
